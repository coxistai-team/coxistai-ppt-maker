"""
AI Content Generation and PowerPoint Creation Module
Handles AI-powered content generation and PowerPoint file creation
"""

import os
import requests
import json
import random
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import logging

logger = logging.getLogger(__name__)

# Configuration
OPENROUTER_API_URL = "https://openrouter.ai/api/v1/chat/completions"
DEFAULT_MODEL = "anthropic/claude-3.5-sonnet"

try:
    from PIL import Image, ImageDraw
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    logger.warning("PIL not available. Image fitting will be limited.")

class SimpleImageGenerator:
    def __init__(self):
        self.api_key = os.getenv("UNSPLASH_API_KEY")
        self.headers = {'Authorization': f'Client-ID {self.api_key}'} if self.api_key else {}
        self.request_count = 0  

    def generate_images(self, topic, num_slides=5):  
        """Get ALL images in ONE request - compatible with your existing code"""
        self.request_count = 0
        try:
            if not self.api_key:
                logger.warning("No Unsplash API key provided. Using placeholders.")
                return self._create_placeholders(topic, num_slides)
            
            # SINGLE API REQUEST
            self.request_count += 1
            response = requests.get(
                "https://api.unsplash.com/search/photos",
                headers=self.headers,
                params={
                    'query': topic,
                    'per_page': num_slides,
                    'orientation': 'landscape'
                },
                timeout=20
            )
            response.raise_for_status()

            os.makedirs("presentation_images", exist_ok=True)
            images = []
            
            # Import S3 service for R2 upload
            try:
                from modules.s3_service import get_s3_service
                s3_service = get_s3_service()
                s3_available = s3_service.is_available()
            except ImportError:
                s3_available = False
            
            for photo in response.json()['results'][:num_slides]:
                img_id = photo['id']
                filename = f"{topic[:20]}_{img_id}.jpg"
                filepath = os.path.join("presentation_images", filename)
                
                if not os.path.exists(filepath):
                    with open(filepath, 'wb') as f:
                        f.write(requests.get(photo['urls']['regular']).content)
                
                # Upload to R2 if available
                s3_url = None
                if s3_available:
                    try:
                        s3_url = s3_service.upload_file(filepath, f"images_{topic}", filename, 'images')
                        logger.info(f"Uploaded image to R2: {s3_url}")
                    except Exception as e:
                        logger.error(f"Failed to upload image to R2: {e}")
                
                images.append({
                    'filepath': filepath,
                    's3_url': s3_url,
                    'photographer': photo['user']['name']
                })
            
            logger.info(f"Generated {len(images)} images in 1 request")
            return images
            
        except Exception as e:
            logger.error(f"Error generating images: {str(e)}")
            return self._create_placeholders(topic, num_slides)

    def _create_placeholders(self, topic, num_slides):
        """Fallback image generation"""
        return [{
            'filepath': None,
            'photographer': "Placeholder"
        } for _ in range(num_slides)]

def get_color_theme():
    """Get a professional color theme"""
    themes = [
        {
            'name': 'Professional Blue',
            'primary': RGBColor(31, 73, 125),
            'secondary': RGBColor(68, 114, 196),
            'accent': RGBColor(255, 193, 7),
            'background': RGBColor(255, 255, 255),
            'card_bg': RGBColor(248, 249, 250),
            'header_bg': RGBColor(31, 73, 125),
            'border': RGBColor(222, 226, 230),
            'text': RGBColor(33, 37, 41),
            'light_text': RGBColor(108, 117, 125),
            'font_primary': 'Calibri',
            'font_secondary': 'Calibri'
        },
        {
            'name': 'Modern Dark',
            'primary': RGBColor(52, 58, 64),
            'secondary': RGBColor(73, 80, 87),
            'accent': RGBColor(0, 123, 255),
            'background': RGBColor(248, 249, 250),
            'card_bg': RGBColor(255, 255, 255),
            'header_bg': RGBColor(52, 58, 64),
            'border': RGBColor(222, 226, 230),
            'text': RGBColor(33, 37, 41),
            'light_text': RGBColor(108, 117, 125),
            'font_primary': 'Segoe UI',
            'font_secondary': 'Segoe UI'
        },
        {
            'name': 'Elegant Purple',
            'primary': RGBColor(102, 51, 153),
            'secondary': RGBColor(147, 112, 219),
            'accent': RGBColor(255, 215, 0),
            'background': RGBColor(255, 255, 255),
            'card_bg': RGBColor(248, 249, 250),
            'header_bg': RGBColor(102, 51, 153),
            'border': RGBColor(222, 226, 230),
            'text': RGBColor(33, 37, 41),
            'light_text': RGBColor(108, 117, 125),
            'font_primary': 'Georgia',
            'font_secondary': 'Georgia'
        }
    ]
    return random.choice(themes)

def fit_image_to_shape(image_path, target_width, target_height):
    """Fit image to target dimensions while maintaining aspect ratio"""
    if not PIL_AVAILABLE or not image_path or not os.path.exists(image_path):
        return None
    
    try:
        with Image.open(image_path) as img:
            # Convert to RGB if necessary
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Calculate aspect ratios
            img_aspect = img.width / img.height
            target_aspect = target_width / target_height
            
            if img_aspect > target_aspect:
                # Image is wider, crop width
                new_width = int(target_height * img_aspect)
                new_height = target_height
                left = (new_width - target_width) // 2
                img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                img = img.crop((left, 0, left + target_width, target_height))
            else:
                # Image is taller, crop height
                new_width = target_width
                new_height = int(target_width / img_aspect)
                top = (new_height - target_height) // 2
                img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                img = img.crop((0, top, target_width, top + target_height))
            
            # Save to temporary file
            temp_path = f"{image_path}_fitted.jpg"
            img.save(temp_path, 'JPEG', quality=85)
            return temp_path
            
    except Exception as e:
        logger.error(f"Error fitting image: {e}")
        return None

def create_thank_you_slide(prs, COLORS, topic):
    """Create a professional thank you slide"""
    thank_you_slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = thank_you_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['background']
    background.line.fill.background()

    # Accent shape
    accent_shape = thank_you_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        Inches(4), Inches(7.5)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = COLORS['primary']
    accent_shape.line.fill.background()

    # Overlay
    overlay = thank_you_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(3.5), Inches(0), 
        Inches(1), Inches(7.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLORS['secondary']
    overlay.fill.transparency = 0.3
    overlay.line.fill.background()

    # Thank you text
    thank_you_box = thank_you_slide.shapes.add_textbox(
        Inches(4.8), Inches(2.5), Inches(8), Inches(1.5)
    )
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "THANK YOU"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.color.rgb = COLORS['text']
    thank_you_para.font.size = Pt(54)
    thank_you_para.font.bold = True
    thank_you_para.font.name = COLORS['font_primary']
    thank_you_para.alignment = PP_ALIGN.LEFT

    # Subtitle
    subtitle_box = thank_you_slide.shapes.add_textbox(
        Inches(4.8), Inches(4.2), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"Questions & Discussion"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.color.rgb = COLORS['light_text']
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.name = COLORS['font_secondary']
    subtitle_para.alignment = PP_ALIGN.LEFT

    # Topic
    topic_box = thank_you_slide.shapes.add_textbox(
        Inches(4.8), Inches(5.5), Inches(8), Inches(0.8)
    )
    topic_frame = topic_box.text_frame
    topic_frame.text = f"Presentation on: {topic}"
    topic_para = topic_frame.paragraphs[0]
    topic_para.font.color.rgb = COLORS['light_text']
    topic_para.font.size = Pt(16)
    topic_para.font.name = COLORS['font_secondary']
    topic_para.alignment = PP_ALIGN.LEFT

    # Accent line
    accent_line = thank_you_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.8), Inches(4), 
        Inches(3), Inches(0.08)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = COLORS['accent']
    accent_line.line.fill.background()

def generate_ai_content(topic: str, num_slides: int, api_key: str) -> Optional[List[Dict[str, Any]]]:
    """
    Generate AI-powered presentation content using OpenRouter API
    
    Args:
        topic: The presentation topic
        num_slides: Number of slides to generate
        api_key: OpenRouter API key
        
    Returns:
        List of slide content dictionaries or None if failed
    """
    if not api_key:
        logger.error("No OpenRouter API key provided")
        return None
    
    try:
        # Create a comprehensive prompt for presentation generation
        prompt = f"""
        Create a professional presentation about "{topic}" with {num_slides} slides.
        
        For each slide, provide:
        1. A clear, engaging title
        2. Key points or bullet points (3-5 points per slide)
        3. A brief description of what should be included
        
        IMPORTANT: You must respond with ONLY a valid JSON array. Do not include any other text, explanations, or markdown formatting.
        
        Format the response as a JSON array with each slide containing:
        - title: The slide title (string)
        - content: Main content points as an array of strings
        - description: Brief description of the slide's purpose (string)
        
        Example format:
        [
          {{
            "title": "Introduction",
            "content": ["Point 1", "Point 2", "Point 3"],
            "description": "Overview of the topic"
          }},
          {{
            "title": "Key Concepts",
            "content": ["Concept 1", "Concept 2", "Concept 3"],
            "description": "Main concepts to understand"
          }}
        ]
        
        Make the content informative, engaging, and well-structured.
        Focus on the most important aspects of {topic}.
        """
        
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
            "HTTP-Referer": "https://coxist-ai.com",
            "X-Title": "AI Presentation Generator"
        }
        
        data = {
            "model": DEFAULT_MODEL,
            "messages": [
                {
                    "role": "system",
                    "content": "You are an expert presentation creator. Generate clear, engaging, and professional presentation content."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            "temperature": 0.7,
            "max_tokens": 2000
        }
        
        logger.info(f"Generating AI content for topic: {topic}")
        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data, timeout=30)
        
        if response.status_code == 200:
            result = response.json()
            content = result.get('choices', [{}])[0].get('message', {}).get('content', '')
            
            try:
                # Parse the JSON response
                slides_data = json.loads(content)
                if isinstance(slides_data, list) and len(slides_data) > 0:
                    logger.info(f"Successfully generated {len(slides_data)} slides")
                    return slides_data
                else:
                    logger.error("Invalid JSON structure in AI response")
                    return create_fallback_slides(topic, num_slides, content)
            except json.JSONDecodeError as e:
                logger.error(f"Failed to parse AI response as JSON: {e}")
                return create_fallback_slides(topic, num_slides, content)
        else:
            logger.error(f"API request failed with status {response.status_code}")
            return create_fallback_slides(topic, num_slides)
            
    except Exception as e:
        logger.error(f"Error generating AI content: {str(e)}")
        return create_fallback_slides(topic, num_slides)

def create_fallback_slides(topic: str, num_slides: int, ai_content: str = "") -> List[Dict[str, Any]]:
    """
    Create fallback slides when AI generation fails
    
    Args:
        topic: The presentation topic
        num_slides: Number of slides to create
        ai_content: Any AI content that was generated (for debugging)
        
    Returns:
        List of basic slide content dictionaries
    """
    logger.info(f"Creating fallback slides for topic: {topic}")
    
    slides = []
    
    # Title slide
    slides.append({
        "title": f"Introduction to {topic}",
        "content": [
            f"Welcome to our presentation on {topic}",
            "This presentation will cover key aspects and insights",
            "Let's explore this fascinating topic together"
        ],
        "description": "Introduction and overview"
    })
    
    # Content slides
    for i in range(1, min(num_slides, 6)):
        if i == 1:
            slides.append({
                "title": f"Key Concepts of {topic}",
                "content": [
                    "Understanding the fundamental principles",
                    "Exploring core methodologies and approaches",
                    "Identifying important trends and developments"
                ],
                "description": "Core concepts and principles"
            })
        elif i == 2:
            slides.append({
                "title": f"Applications and Use Cases",
                "content": [
                    "Real-world applications and implementations",
                    "Industry-specific use cases and examples",
                    "Practical benefits and advantages"
                ],
                "description": "Practical applications and examples"
            })
        elif i == 3:
            slides.append({
                "title": f"Challenges and Considerations",
                "content": [
                    "Common challenges and obstacles",
                    "Important considerations and limitations",
                    "Risk factors and mitigation strategies"
                ],
                "description": "Challenges and considerations"
            })
        elif i == 4:
            slides.append({
                "title": f"Future Trends and Opportunities",
                "content": [
                    "Emerging trends and developments",
                    "Future opportunities and potential",
                    "Recommendations for next steps"
                ],
                "description": "Future outlook and opportunities"
            })
        else:
            slides.append({
                "title": f"Additional Insights on {topic}",
                "content": [
                    "Further exploration of key aspects",
                    "Additional perspectives and viewpoints",
                    "Supporting evidence and data"
                ],
                "description": "Additional insights and information"
            })
    
    # Summary slide
    slides.append({
        "title": "Summary and Conclusion",
        "content": [
            f"Key takeaways from our discussion of {topic}",
            "Important points to remember",
            "Thank you for your attention"
        ],
        "description": "Summary and conclusion"
    })
    
    return slides

def create_powerpoint(slides_data: List[Dict[str, Any]], topic: str) -> Optional[str]:
    """
    Create a professional PowerPoint presentation with enhanced design
    
    Args:
        slides_data: List of slide content dictionaries
        topic: The presentation topic
        
    Returns:
        Path to the created PowerPoint file or None if failed
    """
    try:
        image_generator = SimpleImageGenerator()
        COLORS = get_color_theme()
        
        logger.info(f"Creating presentation with {COLORS['name']} theme...")
        logger.info(f"Typography: {COLORS['font_primary']} (Headers) + {COLORS['font_secondary']} (Body)")

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Generate images for slides
        logger.info(f"Fetching {len(slides_data)} images for the topic '{topic}'...")
        all_images = image_generator.generate_images(topic, num_slides=len(slides_data))

        slide_images = {}
        for idx, slide_data in enumerate(slides_data):
            if all_images and idx < len(all_images) and all_images[idx]['filepath']:
                slide_images[idx] = all_images[idx]['filepath']
                logger.info(f"Assigned image for slide {idx + 1}: {slide_data['title']}")
            else:
                slide_images[idx] = None
                logger.info(f"Using placeholder for slide {idx + 1}: {slide_data['title']}")

        # TITLE SLIDE
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        background = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), 
            Inches(13.33), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = COLORS['background']
        background.line.fill.background()

        accent_shape = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), 
            Inches(4), Inches(7.5)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = COLORS['primary']
        accent_shape.line.fill.background()

        overlay = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(3.5), Inches(0), 
            Inches(1), Inches(7.5)
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = COLORS['secondary']
        overlay.fill.transparency = 0.3
        overlay.line.fill.background()

        title_box = title_slide.shapes.add_textbox(
            Inches(4.8), Inches(2.5), Inches(8), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = topic.upper()
        title_para = title_frame.paragraphs[0]
        title_para.font.color.rgb = COLORS['text']
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.name = COLORS['font_primary']
        title_para.alignment = PP_ALIGN.LEFT

        subtitle_box = title_slide.shapes.add_textbox(
            Inches(4.8), Inches(4.8), Inches(8), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Professional Presentation • {datetime.now().strftime('%B %Y')}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.color.rgb = COLORS['light_text']
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.name = COLORS['font_secondary']
        subtitle_para.alignment = PP_ALIGN.LEFT

        accent_line = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.8), Inches(4.3), 
            Inches(3), Inches(0.08)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = COLORS['accent']
        accent_line.line.fill.background()

        # CONTENT SLIDES
        for idx, slide_data in enumerate(slides_data, 1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(0), Inches(0), 
                Inches(13.33), Inches(7.5)
            )
            background.fill.solid()
            background.fill.fore_color.rgb = COLORS['background']
            background.line.fill.background()

            content_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(0.4), Inches(0.4), 
                Inches(12.53), Inches(6.7)
            )
            content_card.fill.solid()
            content_card.fill.fore_color.rgb = COLORS['card_bg']
            content_card.line.color.rgb = COLORS['border']
            content_card.line.width = Pt(1)

            header_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(0.4), Inches(0.4), 
                Inches(12.53), Inches(1.1)
            )
            header_bg.fill.solid()
            header_bg.fill.fore_color.rgb = COLORS['header_bg']
            header_bg.line.fill.background()

            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.6), Inches(11.73), Inches(0.7)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_frame.margin_left = Inches(0.2)
            title_frame.margin_right = Inches(0.2)
            title_frame.margin_top = Inches(0.1)
            title_frame.margin_bottom = Inches(0.1)
            
            title_para = title_frame.paragraphs[0]
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.font.size = Pt(24)
            title_para.font.bold = True
            title_para.font.name = COLORS['font_primary']
            title_para.alignment = PP_ALIGN.LEFT

            # Add image if available
            if idx - 1 in slide_images and slide_images[idx - 1]:
                try:
                    image_path = slide_images[idx - 1]
                    if image_path and os.path.exists(image_path):
                        # Fit image to slide
                        fitted_image = fit_image_to_shape(image_path, 400, 300)
                        if fitted_image:
                            slide.shapes.add_picture(
                                fitted_image, 
                                Inches(8), Inches(1.8), 
                                Inches(4.5), Inches(3.4)
                            )
                            # Clean up temporary file
                            if fitted_image != image_path:
                                try:
                                    os.remove(fitted_image)
                                except:
                                    pass
                        else:
                            slide.shapes.add_picture(
                                image_path, 
                                Inches(8), Inches(1.8), 
                                Inches(4.5), Inches(3.4)
                            )
                except Exception as e:
                    logger.error(f"Error adding image to slide {idx}: {e}")

            content_box = slide.shapes.add_textbox(
                Inches(0.9), Inches(1.8), Inches(6.8), Inches(4.5)
            )
            content_frame = content_box.text_frame
            content_frame.clear()
            content_frame.word_wrap = True
            content_frame.margin_left = Inches(0.2)
            content_frame.margin_top = Inches(0.2)
            content_frame.margin_right = Inches(0.2)
            content_frame.margin_bottom = Inches(0.2)

            for i, point in enumerate(slide_data['content'][:4]):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                
                p.text = f"• {point}"
                p.font.size = Pt(18)
                p.font.name = COLORS['font_secondary']
                p.font.color.rgb = COLORS['text']
                p.space_after = Pt(12)

        # Add thank you slide
        create_thank_you_slide(prs, COLORS, topic)
        
        # Save the presentation
        filename = f"presentation_{topic.replace(' ', '_').lower()}_{int(datetime.now().timestamp())}.pptx"
        filepath = os.path.join("generated_ppts", filename)
        
        # Ensure directory exists
        os.makedirs("generated_ppts", exist_ok=True)
        
        prs.save(filepath)
        logger.info(f"PowerPoint created successfully: {filepath}")
        return filepath
        
    except Exception as e:
        logger.error(f"Error creating PowerPoint: {str(e)}")
        return None

def enhance_slide_design(slide, slide_data: Dict[str, Any]):
    """
    Enhance slide design with better formatting and styling
    
    Args:
        slide: The slide object to enhance
        slide_data: Data for the slide
    """
    try:
        # Set title formatting
        if slide.shapes.title:
            title = slide.shapes.title
            title.text_frame.paragraphs[0].font.size = Pt(44)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
        
        # Set content formatting
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        paragraph.font.size = Pt(18)
                        paragraph.font.color.rgb = RGBColor(51, 51, 51)
                        
    except Exception as e:
        logger.error(f"Error enhancing slide design: {str(e)}")

def create_enhanced_powerpoint(slides_data: List[Dict[str, Any]], topic: str) -> Optional[str]:
    """
    Create an enhanced PowerPoint presentation with professional design
    
    Args:
        slides_data: List of slide content dictionaries
        topic: The presentation topic
        
    Returns:
        Path to the created PowerPoint file or None if failed
    """
    try:
        # Use the enhanced create_powerpoint function
        return create_powerpoint(slides_data, topic)
        
    except Exception as e:
        logger.error(f"Error creating enhanced PowerPoint: {str(e)}")
        return None

def create_enhanced_slide(prs, slide_data: Dict[str, Any], slide_index: int):
    """
    Create an enhanced slide with professional design
    
    Args:
        prs: The presentation object
        slide_data: Data for the slide
        slide_index: Index of the slide
    """
    try:
        # Create slide with blank layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.33), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        background.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11.33), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', f'Slide {slide_index + 1}')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 73, 125)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add content
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.33), Inches(4.5)
        )
        content_frame = content_box.text_frame
        content_frame.clear()
        
        content = slide_data.get('content', [])
        for i, point in enumerate(content):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.space_after = Pt(12)
            
    except Exception as e:
        logger.error(f"Error creating enhanced slide: {str(e)}")

def create_basic_slide(prs, slide_data: Dict[str, Any], slide_index: int, topic: str):
    """
    Create a basic slide with simple layout
    
    Args:
        prs: The presentation object
        slide_data: Data for the slide
        slide_index: Index of the slide
        topic: The presentation topic
    """
    try:
        if slide_index == 0:
            # Title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = slide_data.get('title', topic)
            
            if slide.shapes.placeholders:
                subtitle = slide.shapes.placeholders[1]
                subtitle.text = f"AI Generated Presentation\n{topic}"
        else:
            # Content slide
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = slide_data.get('title', f'Slide {slide_index + 1}')
            
            content = slide_data.get('content', [])
            if content and slide.shapes.placeholders:
                content_placeholder = slide.shapes.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                for j, point in enumerate(content):
                    if j == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    
    except Exception as e:
        logger.error(f"Error creating basic slide: {str(e)}") 

def create_powerpoint_from_rich_slides(slides_data: List[Dict[str, Any]], topic: str) -> Optional[str]:
    """
    Create a PowerPoint presentation from rich frontend slide data
    
    Args:
        slides_data: List of rich slide data with elements and background
        topic: The presentation topic
        
    Returns:
        Path to the created PowerPoint file or None if failed
    """
    try:
        image_generator = SimpleImageGenerator()
        COLORS = get_color_theme()
        
        logger.info(f"Creating presentation from rich slides with {COLORS['name']} theme...")
        logger.info(f"Typography: {COLORS['font_primary']} (Headers) + {COLORS['font_secondary']} (Body)")

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Generate images for slides
        logger.info(f"Fetching {len(slides_data)} images for the topic '{topic}'...")
        all_images = image_generator.generate_images(topic, num_slides=len(slides_data))

        slide_images = {}
        for idx, slide_data in enumerate(slides_data):
            if all_images and idx < len(all_images) and all_images[idx]['filepath']:
                slide_images[idx] = all_images[idx]['filepath']
                logger.info(f"Assigned image for slide {idx + 1}")
            else:
                slide_images[idx] = None
                logger.info(f"Using placeholder for slide {idx + 1}")

        # TITLE SLIDE
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        background = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), 
            Inches(13.33), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = COLORS['background']
        background.line.fill.background()

        accent_shape = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), 
            Inches(4), Inches(7.5)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = COLORS['primary']
        accent_shape.line.fill.background()

        overlay = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(3.5), Inches(0), 
            Inches(1), Inches(7.5)
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = COLORS['secondary']
        overlay.fill.transparency = 0.3
        overlay.line.fill.background()

        title_box = title_slide.shapes.add_textbox(
            Inches(4.8), Inches(2.5), Inches(8), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = topic.upper()
        title_para = title_frame.paragraphs[0]
        title_para.font.color.rgb = COLORS['text']
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.name = COLORS['font_primary']
        title_para.alignment = PP_ALIGN.LEFT

        subtitle_box = title_slide.shapes.add_textbox(
            Inches(4.8), Inches(4.8), Inches(8), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Professional Presentation • {datetime.now().strftime('%B %Y')}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.color.rgb = COLORS['light_text']
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.name = COLORS['font_secondary']
        subtitle_para.alignment = PP_ALIGN.LEFT

        accent_line = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.8), Inches(4.3), 
            Inches(3), Inches(0.08)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = COLORS['accent']
        accent_line.line.fill.background()

        # CONTENT SLIDES
        for idx, slide_data in enumerate(slides_data, 1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(0), Inches(0), 
                Inches(13.33), Inches(7.5)
            )
            background.fill.solid()
            background.fill.fore_color.rgb = COLORS['background']
            background.line.fill.background()

            content_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(0.4), Inches(0.4), 
                Inches(12.53), Inches(6.7)
            )
            content_card.fill.solid()
            content_card.fill.fore_color.rgb = COLORS['card_bg']
            content_card.line.color.rgb = COLORS['border']
            content_card.line.width = Pt(1)

            header_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(0.4), Inches(0.4), 
                Inches(12.53), Inches(1.1)
            )
            header_bg.fill.solid()
            header_bg.fill.fore_color.rgb = COLORS['header_bg']
            header_bg.line.fill.background()

            # Process rich slide elements
            elements = slide_data.get('elements', [])
            title_text = ""
            content_items = []
            
            for element in elements:
                if element.get('type') == 'title' and element.get('content'):
                    title_text = element['content']
                elif element.get('type') == 'text' and element.get('content'):
                    content_items.append(element['content'])
                elif element.get('type') == 'bullet_list' and element.get('items'):
                    content_items.extend(element['items'])

            # Add title
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.6), Inches(11.73), Inches(0.7)
            )
            title_frame = title_box.text_frame
            title_frame.text = title_text or f"Slide {idx}"
            title_frame.margin_left = Inches(0.2)
            title_frame.margin_right = Inches(0.2)
            title_frame.margin_top = Inches(0.1)
            title_frame.margin_bottom = Inches(0.1)
            
            title_para = title_frame.paragraphs[0]
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.font.size = Pt(24)
            title_para.font.bold = True
            title_para.font.name = COLORS['font_primary']
            title_para.alignment = PP_ALIGN.LEFT

            # Add image if available
            if idx - 1 in slide_images and slide_images[idx - 1]:
                try:
                    image_path = slide_images[idx - 1]
                    if image_path and os.path.exists(image_path):
                        # Fit image to slide
                        fitted_image = fit_image_to_shape(image_path, 400, 300)
                        if fitted_image:
                            slide.shapes.add_picture(
                                fitted_image, 
                                Inches(8), Inches(1.8), 
                                Inches(4.5), Inches(3.4)
                            )
                            # Clean up temporary file
                            if fitted_image != image_path:
                                try:
                                    os.remove(fitted_image)
                                except:
                                    pass
                        else:
                            slide.shapes.add_picture(
                                image_path, 
                                Inches(8), Inches(1.8), 
                                Inches(4.5), Inches(3.4)
                            )
                except Exception as e:
                    logger.error(f"Error adding image to slide {idx}: {e}")

            # Add content
            content_box = slide.shapes.add_textbox(
                Inches(0.9), Inches(1.8), Inches(6.8), Inches(4.5)
            )
            content_frame = content_box.text_frame
            content_frame.clear()
            content_frame.word_wrap = True
            content_frame.margin_left = Inches(0.2)
            content_frame.margin_top = Inches(0.2)
            content_frame.margin_right = Inches(0.2)
            content_frame.margin_bottom = Inches(0.2)

            for i, item in enumerate(content_items[:4]):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.font.name = COLORS['font_secondary']
                p.font.color.rgb = COLORS['text']
                p.space_after = Pt(12)

        # Add thank you slide
        create_thank_you_slide(prs, COLORS, topic)
        
        # Save the presentation
        filename = f"presentation_{topic.replace(' ', '_').lower()}_{int(datetime.now().timestamp())}.pptx"
        filepath = os.path.join("generated_ppts", filename)
        
        # Ensure directory exists
        os.makedirs("generated_ppts", exist_ok=True)
        
        prs.save(filepath)
        logger.info(f"PowerPoint created successfully from rich slides: {filepath}")
        return filepath
        
    except Exception as e:
        logger.error(f"Error creating PowerPoint from rich slides: {str(e)}")
        return None 